"""生成 青野食研所 企业团餐合作方案 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# colors
TEA_GREEN = RGBColor(0x7B, 0xA0, 0x5B)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
DARK = RGBColor(0x3D, 0x3D, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREEN = RGBColor(0xE8, 0xF0, 0xE0)
MID_GREEN = RGBColor(0x9B, 0xBD, 0x7A)
GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT_GRAY = RGBColor(0xF0, 0xEE, 0xEB)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, color, alpha=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_text_box(slide, left, top, w, h, text, font_size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_para(text_frame, text, font_size=16, color=DARK, bold=False, align=PP_ALIGN.LEFT, space_before=Pt(6), font_name='微软雅黑'):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p

def page_number(slide, num):
    add_text_box(slide, W - Inches(1), H - Inches(0.5), Inches(0.8), Inches(0.4),
                 str(num), font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)

# ─── Slide 1: Cover ───
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, CREAM)
# left green block
add_shape(slide, 0, 0, Inches(5.8), H, TEA_GREEN)
# title on green
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(4.5), Inches(1.2),
             '青野食研所', font_size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(3.0), Inches(4.5), Inches(0.6),
             'QINGYE FOOD LAB', font_size=20, color=RGBColor(0xCD, 0xE0, 0xC0))
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(4.5), Inches(0.8),
             '企业团餐合作方案', font_size=36, color=WHITE, bold=True)
# right side text
add_text_box(slide, Inches(6.5), Inches(3.0), Inches(5.5), Inches(0.5),
             '健康 · 新鲜 · 便捷', font_size=24, color=TEA_GREEN, bold=True, align=PP_ALIGN.RIGHT)
add_text_box(slide, Inches(6.5), Inches(3.7), Inches(5.5), Inches(1.5),
             '为您的团队提供高品质午餐配送服务\n让每一位员工吃好，工作更好',
             font_size=16, color=DARK, align=PP_ALIGN.RIGHT)
# decorative dot
for i in range(3):
    y = Inches(5.5) + i * Inches(0.4)
    add_shape(slide, Inches(6.5) + i * Inches(0.35), y, Inches(0.12), Inches(0.12), MID_GREEN)

# ─── Slide 2: 关于我们 ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, Inches(0.15), H, TEA_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.7),
             '关于青野', font_size=32, color=TEA_GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.4),
             'ABOUT QINGYE', font_size=12, color=GRAY)

# left column
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '青野食研所成立于 2023 年，专注于健康轻食的研发与配送。'
p.font.size = Pt(15)
p.font.color.rgb = DARK
p.font.name = '微软雅黑'
for t in [
    '核心团队来自餐饮与营养学背景，坚持「食材干净、调味克制、出品稳定」的产品理念。',
    '',
    '📍 位于 CBD 核心商圈，覆盖 3 公里半径即时配送',
    '📦 日均产能 500+ 份，可承接 50-200 人规模的企业团餐',
    '⭐ 已服务 30+ 家企业，好评率 98%',
]:
    add_para(tf, t, font_size=15, color=DARK, space_before=Pt(8))

# right card
card = add_shape(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.8), LIGHT_GREEN)
add_text_box(slide, Inches(7.3), Inches(2.0), Inches(5), Inches(0.5),
             '我们的优势', font_size=22, color=TEA_GREEN, bold=True)
items = [
    ('🥗 新鲜现做', '每日凌晨采购，当天制作，拒绝隔夜'),
    ('🚚 准时送达', '承诺 11:30 前送达，不影响午休'),
    ('📋 菜单轮换', '每周 5 套不同菜单，吃不腻'),
    ('💰 企业专享价', '团餐均价比单点低 20-30%'),
    ('♻️ 环保包装', '可降解餐盒 + 保温袋配送'),
]
txBox2 = slide.shapes.add_textbox(Inches(7.3), Inches(2.7), Inches(4.8), Inches(3.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
first = True
for title, desc in items:
    if first:
        p = tf2.paragraphs[0]
        first = False
    else:
        p = tf2.add_paragraph()
    run = p.add_run()
    run.text = title
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = DARK
    run.font.name = '微软雅黑'
    run2 = p.add_run()
    run2.text = f'\n{desc}'
    run2.font.size = Pt(13)
    run2.font.color.rgb = GRAY
    run2.font.name = '微软雅黑'
    p.space_before = Pt(10)
page_number(slide, 2)

# ─── Slide 3: 服务模式 ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, Inches(0.15), H, TEA_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.7),
             '服务模式', font_size=32, color=TEA_GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.4),
             'HOW IT WORKS', font_size=12, color=GRAY)

steps = [
    ('01', '确认合作', '沟通企业人数、预算、用餐偏好\n签订月度配送协议'),
    ('02', '定制菜单', '根据企业需求搭配每周菜单\n支持员工提前一天选餐'),
    ('03', '每日配送', '中央厨房统一制作\n保温箱配送至企业前台'),
    ('04', '售后反馈', '每月收集用餐反馈\n动态调整菜单'),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.8) + i * Inches(3.1)
    y = Inches(2.0)
    add_shape(slide, x, y, Inches(2.8), Inches(4.5), LIGHT_GREEN)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.8), Inches(0.6),
                 num, font_size=36, color=TEA_GREEN, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(2.2), Inches(0.5),
                 title, font_size=20, color=DARK, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.6), Inches(2.2), Inches(2.5),
                 desc, font_size=14, color=DARK)
page_number(slide, 3)

# ─── Slide 4: 包月套餐对比 ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, Inches(0.15), H, TEA_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
             '包月套餐对比', font_size=32, color=TEA_GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.4),
             'MEAL PLANS', font_size=12, color=GRAY)

plans = [
    ('轻享版', '¥28 / 餐', '适合预算优先\n的基本午餐需求',
     ['主食（三选一）', '饮品 1 份', '无甜品', '每周 3 套菜单轮换'],
     LIGHT_GREEN, TEA_GREEN),
    ('标准版 ★', '¥38 / 餐', '性价比之选\n最受欢迎',
     ['主食 + 沙拉（各选一）', '饮品 1 份', '甜品/水果 1 份', '每周 5 套菜单轮换',
      '支持提前一日选餐'],
     TEA_GREEN, WHITE),
    ('尊享版', '¥48 / 餐', '全品类自由搭配\n的高端体验',
     ['主食 + 沙拉 + 小吃（各选一）', '饮品 1 份', '甜品 + 水果各 1 份',
      '每周 5 套菜单轮换', '支持提前一周选餐', '月度健康报告'],
     LIGHT_GREEN, TEA_GREEN),
]
for i, (name, price, subtitle, features, bg_color, text_color) in enumerate(plans):
    x = Inches(0.8) + i * Inches(4.1)
    y = Inches(2.0)
    card_w = Inches(3.7)
    card_h = Inches(4.8)
    add_shape(slide, x, y, card_w, card_h, bg_color)
    # header
    hdr = add_shape(slide, x, y, card_w, Inches(1.2), TEA_GREEN)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), card_w - Inches(0.4), Inches(0.5),
                 name, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.55), card_w - Inches(0.4), Inches(0.5),
                 price, font_size=28, color=WHITE, bold=True)
    # subtitle
    add_text_box(slide, x + Inches(0.2), y + Inches(1.3), card_w - Inches(0.4), Inches(0.7),
                 subtitle, font_size=13, color=DARK if bg_color == LIGHT_GREEN else TEA_GREEN)
    # features
    txBox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(2.0), card_w - Inches(0.4), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, feat in enumerate(features):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f'✓ {feat}'
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.font.name = '微软雅黑'
        p.space_before = Pt(6)
    if i == 1:
        add_text_box(slide, x + Inches(0.3), y + Inches(4.2), card_w - Inches(0.6), Inches(0.4),
                     '👍 推荐选择', font_size=13, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER)
        add_shape(slide, x + Inches(0.3), y + Inches(4.6), card_w - Inches(0.6), Inches(0.1), WHITE)

page_number(slide, 4)

# ─── Slide 5: 本周菜单示例 ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, Inches(0.15), H, TEA_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
             '本周菜单示例', font_size=32, color=TEA_GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.4),
             'SAMPLE MENU', font_size=12, color=GRAY)

days = ['周一', '周二', '周三', '周四', '周五']
menus = [
    ['烤鸡胸沙拉 + 南瓜浓汤', '金枪鱼波奇饭 + 抹茶拿铁', '经典凯撒沙拉 + 柠檬水'],
    ['照烧鸡腿三明治 + 水果杯', '牛油果拌饭 + 味噌汤', '地中海沙拉 + 气泡水'],
    ['烟熏三文鱼波奇饭', '黑椒牛肉卷饼 + 蔬菜杯', '泰式虾仁沙拉 + 椰子水'],
    ['日式鸡排饭 + 泡菜', '吞拿鱼三明治 + 酸奶', '考伯沙拉 + 蜂蜜柠檬茶'],
    ['香草鸡胸全麦卷', '韩式拌饭 + 海带汤', '鲜虾牛油果沙拉 + 豆浆'],
]
for i, day in enumerate(days):
    x = Inches(0.6) + i * Inches(2.5)
    y = Inches(1.9)
    add_shape(slide, x, y, Inches(2.3), Inches(5), LIGHT_GREEN)
    add_shape(slide, x, y, Inches(2.3), Inches(0.7), TEA_GREEN)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(2), Inches(0.5),
                 day, font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.9), Inches(2), Inches(3.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, menu in enumerate(menus[i]):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = menu
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.font.name = '微软雅黑'
        p.space_before = Pt(12)
page_number(slide, 5)

# ─── Slide 6: 配送流程 ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, Inches(0.15), H, TEA_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.7),
             '配送流程', font_size=32, color=TEA_GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.4),
             'DELIVERY PROCESS', font_size=12, color=GRAY)

flow = [
    ('📅', '前日 17:00', '员工通过小程序\n完成次日选餐'),
    ('🌅', '当日 07:00', '中央厨房开始\n制作与打包'),
    ('🚚', '当日 10:30', '保温车配送出发'),
    ('✅', '当日 11:30 前', '送达企业前台\n完成签收'),
]
for i, (icon, time, desc) in enumerate(flow):
    x = Inches(0.8) + i * Inches(3.1)
    y = Inches(2.2)
    add_shape(slide, x, y, Inches(2.8), Inches(4), LIGHT_GREEN)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.2), Inches(0.6),
                 icon, font_size=36, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(2.2), Inches(0.5),
                 time, font_size=16, color=TEA_GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.7), Inches(2.2), Inches(2),
                 desc, font_size=14, color=DARK, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        add_text_box(slide, x + Inches(2.6), y + Inches(1.5), Inches(0.6), Inches(0.5),
                     '→', font_size=28, color=TEA_GREEN, bold=True, align=PP_ALIGN.CENTER)
page_number(slide, 6)

# ─── Slide 7: 价格方案 ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, Inches(0.15), H, TEA_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
             '价格方案', font_size=32, color=TEA_GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.4),
             'PRICING', font_size=12, color=GRAY)

pricing = [
    ('试吃体验', '前 3 天', '免费', '让团队先体验\n再决定是否合作'),
    ('月度合作', '20 个工作日', '¥28-48 / 餐', '按选定套餐定价\n月结开票'),
    ('季度合作', '60 个工作日', '¥25-43 / 餐', '享 9 折优惠\n优先排单'),
    ('年度合作', '240 个工作日', '¥23-40 / 餐', '享 85 折优惠\n专属菜单定制'),
]
headers = ['方案', '周期', '单价', '说明']
for i, h in enumerate(headers):
    x = Inches(0.8) + i * Inches(3.1)
    add_shape(slide, x, Inches(1.9), Inches(2.9), Inches(0.6), TEA_GREEN)
    add_text_box(slide, x + Inches(0.1), Inches(1.95), Inches(2.7), Inches(0.5),
                 h, font_size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
for row, (a, b, c, d) in enumerate(pricing):
    y = Inches(2.6) + row * Inches(1.1)
    bg = LIGHT_GREEN if row % 2 == 0 else WHITE
    for col, val in enumerate([a, b, c, d]):
        x = Inches(0.8) + col * Inches(3.1)
        add_shape(slide, x, y, Inches(2.9), Inches(0.95), bg)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.2), Inches(2.7), Inches(0.6),
                     val, font_size=14, color=DARK, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(10), Inches(0.5),
             '* 以上价格不含税，餐费月结，可开增值税普通发票', font_size=12, color=GRAY)
page_number(slide, 7)

# ─── Slide 8: 合作流程 ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, Inches(0.15), H, TEA_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.7),
             '合作流程', font_size=32, color=TEA_GREEN, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.4),
             'COOPERATION PROCESS', font_size=12, color=GRAY)

steps = [
    ('STEP 1', '联系我们', '扫码或电话沟通\n了解企业需求'),
    ('STEP 2', '免费试餐', '安排 3 份试吃餐\n体验品质'),
    ('STEP 3', '签约付款', '确认方案后\n签订月度/季度合同'),
    ('STEP 4', '开始配送', '次周起正式配送\n专人跟进服务'),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.8) + i * Inches(3.1)
    y = Inches(2.2)
    add_shape(slide, x, y, Inches(2.8), Inches(4), LIGHT_GREEN)
    add_shape(slide, x + Inches(0.8), y + Inches(0.3), Inches(1.2), Inches(0.4), TEA_GREEN)
    add_text_box(slide, x + Inches(0.8), y + Inches(0.3), Inches(1.2), Inches(0.4),
                 num, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.0), Inches(2.4), Inches(0.5),
                 title, font_size=20, color=TEA_GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.7), Inches(2.4), Inches(2),
                 desc, font_size=14, color=DARK, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(2.6), y + Inches(1.5), Inches(0.6), Inches(0.5),
                     '→', font_size=28, color=TEA_GREEN, bold=True, align=PP_ALIGN.CENTER)
page_number(slide, 8)

# ─── Slide 9: 联系我们 ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CREAM)
add_shape(slide, 0, 0, Inches(5.8), H, TEA_GREEN)
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(4.5), Inches(1),
             '让每一餐\n都值得期待', font_size=40, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(3.3), Inches(4.5), Inches(0.6),
             '期待与您的企业合作', font_size=20, color=RGBColor(0xCD, 0xE0, 0xC0))
# right side
add_text_box(slide, Inches(6.5), Inches(2.0), Inches(5.5), Inches(0.5),
             '联系我们', font_size=28, color=TEA_GREEN, bold=True)
contact_info = [
    '📞 电话：138-0000-8888',
    '📧 邮箱：qingye@foodlab.com',
    '💬 微信：qingye_food',
    '📍 地址：北京市朝阳区建国路 88 号',
]
txBox = slide.shapes.add_textbox(Inches(6.5), Inches(2.8), Inches(5.5), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True
for i, info in enumerate(contact_info):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = info
    p.font.size = Pt(16)
    p.font.color.rgb = DARK
    p.font.name = '微软雅黑'
    p.space_before = Pt(12)

add_text_box(slide, Inches(6.5), Inches(6.0), Inches(5.5), Inches(0.5),
             '扫码添加商务微信，预约免费试餐 →', font_size=14, color=GRAY)
# qr placeholder
qr = add_shape(slide, Inches(10), Inches(5.5), Inches(1.2), Inches(1.2), WHITE)
add_text_box(slide, Inches(10), Inches(5.8), Inches(1.2), Inches(0.6),
             'QR', font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ─── Save ───
output_path = '/home/wuming/me/wo/startup/青野食研所-企业团餐合作方案.pptx'
prs.save(output_path)
print(f'Saved to {output_path}')
